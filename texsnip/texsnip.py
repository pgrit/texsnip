import subprocess
import os
import tempfile
import shutil
import re

def red(txt):
    return "\033[0;31m" + txt + "\033[00m"

def extract_errors(logfile):
    """ Basic parsing of pdflatex log files.

    Tries to find each line that contains an error message and outputs that line and a number
    of following lines.
    """
    with open(logfile) as fp: lines = fp.readlines()

    num_lines = 5
    def concat(i):
        def non_empty():
            for j in range(i, min(i + num_lines, len(lines))):
                if lines[j].strip():
                    yield lines[j]
        return (
            red("==========================================\n") +
            red("LaTeX error:\n") +
            "------------------------------------------\n" +
            "".join([ l for l in non_empty() ]) +
            red("==========================================")
        )

    for i in range(len(lines)):
        if "error:" in lines[i] or "Error:" in lines[i]: yield concat(i)
        elif ".tex:"in lines[i]: yield concat(i)
        elif lines[i].startswith("! "): yield concat(i)

def compile_and_crop(tex, name, intermediate_dir = None):
    """ Compiles the given LaTeX code.

    Args:
        - tex   The file content of the LaTeX file to compile
        - name  Name of the output without the .pdf extension
        - intermediate_dir  Specify an existing directory here and .tex and .log files will be kept there

    Returns:
        True if compilation succeeded, false otherwise. Error messages will be printed in the console.
    """
    if intermediate_dir is not None and os.path.isdir(intermediate_dir):
        temp_folder = None
        temp_dir = os.path.abspath(intermediate_dir)
    else:
        temp_folder = tempfile.TemporaryDirectory()
        temp_dir = temp_folder.name

    with open(os.path.join(temp_dir, f"{name}.tex"), "w") as fp:
        fp.write(tex)

    try:
        subprocess.check_call([
                "pdflatex",
                "-interaction=batchmode",
                f"{name}.tex"
            ], cwd=temp_dir, stdout=subprocess.DEVNULL)
    except subprocess.CalledProcessError:
        logfile = os.path.join(temp_dir, f"{name}.log")
        if not os.path.exists(logfile):
            print(red("Error: pdflatex failed, but no log was written."))
        else:
            print("\n".join([errline for errline in extract_errors(logfile)]))
            print(red(f"Error: pdflatex failed. Syntax error or missing package? "
                "You can view the full log in {logfile}. This path can be changed by specifying an intermediate_dir"))
        return False

    subprocess.check_call(["pdfcrop", f"{name}.pdf"], cwd=temp_dir, stdout=subprocess.DEVNULL)
    shutil.copy(os.path.join(temp_dir, f"{name}-crop.pdf"), f"{name}.pdf")

    if temp_folder is not None:
        temp_folder.cleanup()

    return True

class Snip:
    def __init__(self, name, fontsize_pt, content):
        self._name = str(name)
        self.fontsize_pt = fontsize_pt
        self.content = content

    @property
    def name(self):
        return self._name

    @property
    def fontsize_pt(self):
        return self._fontsize_pt

    @fontsize_pt.setter
    def fontsize_pt(self, value):
        try:
            sz = float(value)
        except:
            raise ValueError("fontsize must be a number (real or integer)")

        if sz <= 0:
            raise ValueError("fontsize must not be negative!")

        self._fontsize_pt = value

    @property
    def content(self):
        return self._content

    @content.setter
    def content(self, value):
        self._content = str(value)

    def generate(self, preamble=r"\usepackage{libertine}", intermediate_dir=None):
        tex_code = "\n".join([
            r"\documentclass{article}",
            r"\pagenumbering{gobble}",
            r"\usepackage{xcolor}",
            r"\usepackage{graphicx}",
            r"\usepackage[utf8]{inputenc}",
            r"\usepackage[T1]{fontenc}",
            r"\usepackage{geometry}",
            r"\geometry{",
            r"    papersize={500cm,500cm},",
            r"    total={500cm,500cm},",
            r"    left=0mm,",
            r"    top=0mm,",
            r"}",
            ""
        ])

        tex_code += preamble
        tex_code += r"\begin{document}"

        # scale the fontsize
        fontcmd = r"{\fontsize{" + f"{self.fontsize_pt}" + "pt}{" + f"{0}" + r"pt}"
        fontcmd += r"\selectfont\raggedright" + "\n"

        # scale math accordingly, using scalebox to also adjust large operators like \sum
        scalecmd = r"}\\scalebox{" + f"{self.fontsize_pt/10}" + "}{"
        fontcmd_regex = r"{\\fontsize{" + f"{self.fontsize_pt}" + "pt}{" + f"{0}" + r"pt}"
        fontcmd_regex += r"\\selectfont\\raggedright" + "\n"
        wrapped_math = re.sub(r"(\$.*\$)", scalecmd + r"\1}" + fontcmd_regex, self.content)

        tex_code += fontcmd + wrapped_math
        tex_code += "\n"

        # tex_code += self.content + "\n"
        tex_code += r"} \end{document}"

        compile_and_crop(tex_code, self.name, intermediate_dir)

    def generate_png(self, preamble=r"\usepackage{libertine}", intermediate_dir=None, dpi=300):
        self.generate(preamble, intermediate_dir)

        from pdf2image import convert_from_path
        convert_from_path(self.name + ".pdf", dpi=dpi, transparent=True, fmt="png", output_file=self.name, single_file=True, output_folder=".")

def pptx_snips(snips, filename="snips.pptx", preamble=r"\usepackage{libertine}", intermediate_dir=None, dpi=300):
    from PyPDF2 import PdfFileReader
    from pptx import Presentation
    from pptx.util import Pt

    # Create a new .pptx with a single slide
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    ypos_pt = 8
    for s in snips:
        s.generate_png(preamble, intermediate_dir, dpi)

        # Compute the pdf file size so we can scale the png to the correct size
        # Not doing so would make fontsizes inconsistent
        box = PdfFileReader(open(s.name + ".pdf", "rb")).getPage(0).mediaBox
        width_pt = box.upperRight[0]
        height_pt = box.upperRight[1]

        # Place the .png with the correct size in points
        slide.shapes.add_picture(s.name + ".png", Pt(8), Pt(ypos_pt), width=Pt(width_pt))

        ypos_pt += height_pt + 8

    prs.save(filename)